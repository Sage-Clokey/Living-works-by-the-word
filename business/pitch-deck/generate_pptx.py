#!/usr/bin/env python3
"""Living Works by the Word — Pitch Deck Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from lxml import etree

BASE = "/mnt/c/Users/SajcS/Desktop/Living works by the word/"
LA   = BASE + "living_architecture_images/"
LOA  = BASE + "images of living age/"
MG   = BASE + "morphogenetic-designer/outputs/frames/"
OUT  = BASE + "business plan pitch prompts/Living_Works_Pitch_Deck_v2.pptx"

IMGS = {
    "cover":    LA  + "ChatGPT Image Feb 28, 2026, 02_33_33 PM.png",
    "block_vs": LA  + "06_block_vs_spiral_manifesto.png",
    "mycelium": LA  + "02_mycelium_network_plan.png",
    "bw_vs_sw": LOA + "ChatGPT Image Mar 6, 2026, 02_32_21 PM.png",
    "platform": LOA + "ChatGPT Image Mar 6, 2026, 02_34_47 PM.png",
    "dome_3d":  LA  + "07_mycelium_dome_3d.png",
    "growth":   MG  + "frame_0070.png",
    "tree":     LA  + "04_branching_tree_structure.png",
    "village":  LA  + "ChatGPT Image Feb 28, 2026, 02_38_41 PM.png",
    "facade":   LA  + "03_living_facade.png",
    "fib_plan": LA  + "01_fibonacci_spiral_house.png",
    "sunflower":LA  + "05_sunflower_room_spiral.png",
    "dome2":    LA  + "ChatGPT Image Feb 28, 2026, 02_37_12 PM.png",
}

# ── Colours ──────────────────────────────────────────────────────────────────
BG    = RGBColor(0x0a, 0x18, 0x0e)
PANEL = RGBColor(0x0f, 0x22, 0x17)
GRN   = RGBColor(0x2d, 0x6a, 0x4f)
TEAL  = RGBColor(0x52, 0xb7, 0x88)
LTGRN = RGBColor(0x95, 0xd5, 0xb2)
WHITE = RGBColor(0xff, 0xff, 0xff)
CREAM = RGBColor(0xd8, 0xf3, 0xdc)
GOLD  = RGBColor(0xe9, 0xc4, 0x6a)
GREY  = RGBColor(0x88, 0xa8, 0x90)

W = Inches(13.33)
H = Inches(7.5)

# ── Setup ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, c):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = c

def add_img(slide, path, l, t, w=None, h=None):
    if w and h: return slide.shapes.add_picture(path, l, t, w, h)
    elif w:     return slide.shapes.add_picture(path, l, t, width=w)
    elif h:     return slide.shapes.add_picture(path, l, t, height=h)
    else:       return slide.shapes.add_picture(path, l, t)

def add_rect(slide, l, t, w, h, fill=None, alpha=100):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha < 100:
            spPr = s._element.find(qn('p:spPr'))
            sf = spPr.find(qn('a:solidFill'))
            if sf is not None:
                sc = sf.find(qn('a:srgbClr'))
                if sc is not None:
                    ae = etree.SubElement(sc, qn('a:alpha'))
                    ae.set('val', str(int(alpha * 1000)))
    else:
        s.fill.background()
    return s

def add_hr(slide, l, t, w, c=TEAL, thickness=Inches(0.045)):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, thickness)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def add_notes(slide, text):
    """Add speaker notes to a slide."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def add_tb(slide, l, t, w, h, lines):
    """lines: list of dicts — text, size, bold, italic, color, align, bullet, sp"""
    shape = slide.shapes.add_textbox(l, t, w, h)
    shape.line.fill.background(); shape.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(ln.get('sp', 4))
        text = ('\u2022  ' + ln['text']) if ln.get('bullet') else ln['text']
        r = p.add_run()
        r.text = text
        r.font.name = 'Calibri Light' if ln.get('light') else 'Calibri'
        r.font.size  = Pt(ln.get('size', 18))
        r.font.bold  = ln.get('bold', False)
        r.font.italic = ln.get('italic', False)
        r.font.color.rgb = ln.get('color', WHITE)
    return shape

# ── Shorthand line-dict constructors ────────────────────────────────────────
def L(text, size=18, bold=False, italic=False, color=WHITE,
      align=PP_ALIGN.LEFT, bullet=False, sp=4, light=False):
    return dict(text=text, size=size, bold=bold, italic=italic,
                color=color, align=align, bullet=bullet, sp=sp, light=light)

def LABEL(n, title): return L(f"{n:02d}  {title}", size=11, bold=True, color=TEAL, sp=0)
def TITLE(t, size=38): return L(t, size=size, bold=True, color=WHITE, light=True, sp=2)
def SUB(t, size=20): return L(t, size=size, color=CREAM, sp=6)
def BUL(t, size=17): return L(t, size=size, bullet=True, color=WHITE, sp=5)
def QTE(t, size=17): return L(t, size=size, italic=True, color=GOLD, sp=8)
def HEAD(t, size=15): return L(t, size=15, bold=True, color=LTGRN, sp=10)
def CAP(t): return L(t, size=11, italic=True, color=GREY, sp=2)
def SPC(): return L('', size=5, sp=0)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
def slide_01_cover():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["cover"], 0, 0, W, H)
    # dark gradient strip at bottom
    add_rect(sl, 0, Inches(3.8), W, Inches(3.7), BG, alpha=82)
    # top thin strip
    add_rect(sl, 0, 0, W, Inches(0.5), BG, alpha=75)
    # teal left accent bar
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(1.2), Inches(4.1), Inches(11), Inches(1.6), [
        L('LIVING WORKS', size=70, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, light=True, sp=0),
    ])
    add_hr(sl, Inches(5.2), Inches(5.7), Inches(3.0), TEAL)
    add_tb(sl, Inches(1.2), Inches(5.8), Inches(11), Inches(0.8), [
        L('From the Age of Blocks to the Living Age',
          size=26, color=LTGRN, align=PP_ALIGN.CENTER, light=True),
    ])
    add_tb(sl, Inches(1.2), Inches(6.65), Inches(11), Inches(0.5), [
        L('Living Works by the Word  ·  March 2026',
          size=14, color=GREY, align=PP_ALIGN.CENTER),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
def slide_02_problem():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["block_vs"], Inches(6.2), 0, Inches(7.13), H)
    add_rect(sl, 0, 0, Inches(6.5), H, PANEL)
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(0.4), Inches(0.3), Inches(5.8), Inches(0.35), [LABEL(1, 'THE PROBLEM')])
    add_tb(sl, Inches(0.4), Inches(0.7), Inches(5.8), Inches(1.4), [
        TITLE('Civilisation\nBuilt Against Life', size=40),
    ])
    add_hr(sl, Inches(0.4), Inches(2.15), Inches(2.2), TEAL)
    add_tb(sl, Inches(0.4), Inches(2.35), Inches(5.7), Inches(4.8), [
        SUB('Modern infrastructure is built on rigid mechanical systems that:', size=17),
        SPC(),
        BUL('Require massive energy to manufacture and replace'),
        BUL('Cannot adapt to changing environments'),
        BUL('Degrade over time and generate waste'),
        BUL('Separate human design entirely from natural systems'),
        SPC(), SPC(),
        L('Construction · Manufacturing · Infrastructure',
          size=14, italic=True, color=GREY),
        L('All still anchored to industrial-age engineering.',
          size=14, italic=True, color=GREY, sp=2),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE OPPORTUNITY
# ════════════════════════════════════════════════════════════════════════════
def slide_03_opportunity():
    sl = new_slide(); set_bg(sl, BG)
    # Mycelium network centred-right
    add_img(sl, IMGS["mycelium"], Inches(4.8), Inches(0.4), Inches(8.2), Inches(7.1))
    # Left dark panel
    add_rect(sl, 0, 0, Inches(5.4), H, BG, alpha=97)

    add_tb(sl, Inches(0.4), Inches(0.3), Inches(4.8), Inches(0.35), [LABEL(2, 'THE OPPORTUNITY')])
    add_tb(sl, Inches(0.4), Inches(0.7), Inches(4.8), Inches(1.5), [
        TITLE('Living Systems\nAlready Solve This', size=36),
    ])
    add_hr(sl, Inches(0.4), Inches(2.25), Inches(2.2), TEAL)
    add_tb(sl, Inches(0.4), Inches(2.45), Inches(4.8), Inches(4.7), [
        HEAD('Nature\'s solutions:'),
        SPC(),
        BUL('Trees — self-assemble load-bearing structures'),
        BUL('Coral — builds mineral architecture over centuries'),
        BUL('Fungal networks — distribute resources with optimal efficiency'),
        BUL('Cells — organise into tissues, organs, and organisms'),
        SPC(), SPC(),
        QTE('"Our tools only let us analyse biology.\nNot design with it."'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE VISION
# ════════════════════════════════════════════════════════════════════════════
def slide_04_vision():
    sl = new_slide(); set_bg(sl, BG)
    # Block World vs Spiral World infographic on left
    add_img(sl, IMGS["bw_vs_sw"], 0, 0, Inches(7.3), H)
    # Right panel
    add_rect(sl, Inches(7.3), 0, Inches(6.03), H, PANEL)
    add_rect(sl, W - Inches(0.07), 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(7.65), Inches(0.3), Inches(5.4), Inches(0.35), [LABEL(3, 'THE VISION')])
    add_tb(sl, Inches(7.65), Inches(0.7), Inches(5.4), Inches(1.2), [
        TITLE('The Living Age', size=42),
    ])
    add_hr(sl, Inches(7.65), Inches(1.95), Inches(2.2), TEAL)
    add_tb(sl, Inches(7.65), Inches(2.15), Inches(5.3), Inches(5.0), [
        QTE('"Enable humanity to design with\nliving systems instead of against them."', size=19),
        SPC(), SPC(),
        HEAD('A new design paradigm:'),
        SPC(),
        BUL('Structures grown, not manufactured'),
        BUL('Materials that absorb carbon and self-repair'),
        BUL('Ecosystems designed for regeneration'),
        BUL('Cities that work with biology, not against it'),
        SPC(), SPC(),
        L('Industrial civilisation built with blocks.',
          size=14, italic=True, color=GREY),
        L('The next civilisation will build with living systems.',
          size=14, italic=True, color=GREY, sp=2),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE PLATFORM (full-bleed infographic)
# ════════════════════════════════════════════════════════════════════════════
def slide_05_platform():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["platform"], 0, 0, W, H)
    # Top label strip
    add_rect(sl, 0, 0, W, Inches(0.5), BG, alpha=80)
    add_tb(sl, Inches(0.4), Inches(0.07), Inches(10), Inches(0.38), [
        LABEL(4, 'THE SOLUTION  ·  Living Works Biological Design Platform'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MORPHOGENESIS ENGINE
# ════════════════════════════════════════════════════════════════════════════
def slide_06_morpho():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["dome_3d"], Inches(7.6), 0, Inches(5.73), H)
    add_rect(sl, 0, 0, Inches(7.8), H, PANEL)
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(0.4), Inches(0.3), Inches(7.1), Inches(0.35), [LABEL(5, 'CORE TECHNOLOGY')])
    add_tb(sl, Inches(0.4), Inches(0.7), Inches(7.1), Inches(1.3), [
        TITLE('Morphogenesis Engine', size=38),
    ])
    add_hr(sl, Inches(0.4), Inches(2.05), Inches(2.5), TEAL)
    add_tb(sl, Inches(0.4), Inches(2.2), Inches(7.1), Inches(5.0), [
        SUB('3D simulation of biological growth — predicting how living structures\ndevelop before committing to biological processes.', size=17),
        SPC(), SPC(),
        HEAD('Capabilities:'),
        SPC(),
        BUL('Branching systems — trees, coral, vasculature, fungal networks'),
        BUL('Nutrient diffusion and resource gradient modelling'),
        BUL('Agent-based cellular growth rules'),
        BUL('Biomechanical forces and structural load feedback'),
        SPC(),
        HEAD('Outputs:'),
        SPC(),
        BUL('3D growth models  ·  Gene pathway suggestions  ·  Growth predictions'),
    ])
    add_tb(sl, Inches(7.7), H - Inches(0.45), Inches(5.4), Inches(0.4), [
        CAP('Mycelium Dome — 3D Living Architecture'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI DESIGN INTERFACE
# ════════════════════════════════════════════════════════════════════════════
def slide_07_ai():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["growth"], 0, 0, Inches(5.8), H)
    add_rect(sl, Inches(5.6), 0, Inches(7.73), H, PANEL)
    add_rect(sl, W - Inches(0.07), 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(6.0), Inches(0.3), Inches(7.0), Inches(0.35), [LABEL(6, 'CORE TECHNOLOGY')])
    add_tb(sl, Inches(6.0), Inches(0.7), Inches(7.0), Inches(1.4), [
        TITLE('AI Biological\nDesign Interface', size=36),
    ])
    add_hr(sl, Inches(6.0), Inches(2.15), Inches(2.5), TEAL)
    add_tb(sl, Inches(6.0), Inches(2.35), Inches(6.9), Inches(1.1), [
        SUB('Natural language → biological design.\nDescribe what you want to grow. The system generates the biology.', size=17),
    ])
    # Prompt example box
    add_rect(sl, Inches(6.0), Inches(3.55), Inches(6.9), Inches(1.05), GRN, alpha=55)
    add_tb(sl, Inches(6.15), Inches(3.65), Inches(6.6), Inches(0.85), [
        L('"Design a branching structure optimised for wind resistance\nand water collection."',
          size=15, italic=True, color=CREAM, align=PP_ALIGN.CENTER),
    ])
    add_tb(sl, Inches(6.0), Inches(4.75), Inches(6.9), Inches(2.5), [
        HEAD('System outputs:'),
        SPC(),
        BUL('3D growth models'),
        BUL('Gene pathway suggestions'),
        BUL('Structural simulations'),
        BUL('Biological design specification ready for lab testing'),
    ])
    add_tb(sl, Inches(0.1), H - Inches(0.45), Inches(5.4), Inches(0.4), [
        CAP('Morphogenetic growth simulation — Living Works engine'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BIOCAD PLATFORM
# ════════════════════════════════════════════════════════════════════════════
def slide_08_biocad():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["tree"], Inches(7.5), 0, Inches(5.83), H)
    add_rect(sl, 0, 0, Inches(7.7), H, PANEL)
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(0.4), Inches(0.3), Inches(7.0), Inches(0.35), [LABEL(7, 'CORE TECHNOLOGY')])
    add_tb(sl, Inches(0.4), Inches(0.7), Inches(7.0), Inches(1.2), [
        TITLE('BioCAD Platform', size=38),
    ])
    add_hr(sl, Inches(0.4), Inches(1.95), Inches(2.5), TEAL)
    add_tb(sl, Inches(0.4), Inches(2.15), Inches(7.0), Inches(5.1), [
        SUB('Computer-aided design for living systems.\nThe same design precision as AutoCAD — applied to biology.', size=17),
        SPC(), SPC(),
        HEAD('Design Tools:'),
        SPC(),
        BUL('Genome editing planning and visualisation'),
        BUL('Gene network and pathway design'),
        BUL('Tissue architecture modelling'),
        BUL('Ecosystem simulation environments'),
        BUL('Living materials library — mycelium · bacterial cellulose · plant scaffolds'),
        SPC(), SPC(),
        QTE('"Like Autodesk for buildings · Unity for games · but for biology."'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVING ARCHITECTURE (full bleed)
# ════════════════════════════════════════════════════════════════════════════
def slide_09_arch():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["village"], 0, 0, W, H)
    add_rect(sl, 0, 0, W, Inches(1.15), BG, alpha=78)
    add_rect(sl, 0, H - Inches(0.9), W, Inches(0.9), BG, alpha=82)

    add_tb(sl, Inches(0.4), Inches(0.07), Inches(8), Inches(0.38), [LABEL(8, 'APPLICATIONS')])
    add_tb(sl, Inches(0.4), Inches(0.45), Inches(9), Inches(0.7), [
        TITLE('Living Architecture', size=40),
    ])
    add_tb(sl, Inches(0.4), H - Inches(0.82), W - Inches(0.8), Inches(0.7), [
        L('Buildings grown from mycelium, engineered plants, and root networks  ·  '
          'Carbon-absorbing  ·  Self-repairing  ·  Alive',
          size=17, color=CREAM, align=PP_ALIGN.CENTER),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ALL APPLICATIONS
# ════════════════════════════════════════════════════════════════════════════
def slide_10_applications():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["facade"], 0, 0, Inches(4.9), H)
    add_rect(sl, Inches(4.7), 0, Inches(8.63), H, PANEL)
    add_rect(sl, W - Inches(0.07), 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(5.1), Inches(0.3), Inches(7.9), Inches(0.35), [LABEL(9, 'APPLICATIONS')])
    add_tb(sl, Inches(5.1), Inches(0.7), Inches(7.9), Inches(1.1), [
        TITLE('Four Industries Transformed', size=34),
    ])
    add_hr(sl, Inches(5.1), Inches(1.85), Inches(2.5), TEAL)

    apps = [
        ('Living Architecture',
         'Buildings grown from mycelium and engineered plants'),
        ('Carbon-Negative Construction',
         'Structures that grow, absorb CO\u2082, and self-repair'),
        ('Ecosystem Engineering',
         'Landscapes that regenerate soil, purify water, restore biodiversity'),
        ('Biological Manufacturing',
         'Grown textiles, composites, and packaging — no petroleum required'),
    ]
    y = Inches(2.1)
    for title_t, desc_t in apps:
        add_rect(sl, Inches(5.0), y, Inches(8.0), Inches(1.2), GRN, alpha=30)
        add_tb(sl, Inches(5.25), y + Inches(0.12), Inches(7.6), Inches(0.5), [
            L(title_t, size=18, bold=True, color=LTGRN),
        ])
        add_tb(sl, Inches(5.25), y + Inches(0.55), Inches(7.6), Inches(0.5), [
            L(desc_t, size=15, color=CREAM),
        ])
        y += Inches(1.3)
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL
# ════════════════════════════════════════════════════════════════════════════
def slide_11_biz():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["fib_plan"], Inches(6.5), Inches(0.5), Inches(6.5), Inches(6.5))
    add_rect(sl, Inches(6.5), 0, Inches(6.83), H, BG, alpha=65)
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(0.4), Inches(0.3), Inches(9), Inches(0.35), [LABEL(10, 'BUSINESS MODEL')])
    add_tb(sl, Inches(0.4), Inches(0.7), Inches(7), Inches(1.0), [
        TITLE('Revenue Streams', size=38),
    ])
    add_hr(sl, Inches(0.4), Inches(1.75), Inches(2.5), TEAL)

    streams = [
        ('SaaS Platform',
         'Academic $50–200/mo  ·  Startup $500/mo  ·  Enterprise $2,000+/mo',
         'Universities  ·  Biotech  ·  Architects  ·  Materials firms'),
        ('AI Design Services',
         'Custom organism-design projects — project-based pricing',
         'Construction  ·  Materials  ·  Biotech companies'),
        ('IP Licensing',
         'Biological growth algorithms  ·  Engineered organisms  ·  Living material systems',
         'Royalty-based across industries'),
    ]

    y = Inches(2.0)
    for name, pricing, users in streams:
        add_rect(sl, Inches(0.3), y, Inches(8.2), Inches(1.45), GRN, alpha=38)
        add_tb(sl, Inches(0.55), y + Inches(0.1), Inches(7.8), Inches(0.48), [
            L(name, size=18, bold=True, color=GOLD),
        ])
        add_tb(sl, Inches(0.55), y + Inches(0.52), Inches(7.8), Inches(0.38), [
            L(pricing, size=14, color=WHITE),
        ])
        add_tb(sl, Inches(0.55), y + Inches(0.88), Inches(7.8), Inches(0.38), [
            L('Users: ' + users, size=12, color=GREY),
        ])
        y += Inches(1.62)

    add_tb(sl, Inches(0.3), Inches(6.95), Inches(9), Inches(0.45), [
        L('Phase 1 (Yrs 1–2): University pilots  ·  Phase 2 (Yrs 3–5): Enterprise  ·  Phase 3 (Yrs 5–10): Living infrastructure',
          size=12, italic=True, color=GREY),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — STRATEGY
# ════════════════════════════════════════════════════════════════════════════
def slide_12_strategy():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["sunflower"], 0, Inches(0.3), Inches(6.5), Inches(6.5))
    add_rect(sl, 0, 0, Inches(6.5), H, BG, alpha=45)
    add_rect(sl, Inches(6.3), 0, Inches(7.03), H, PANEL)
    add_rect(sl, W - Inches(0.07), 0, Inches(0.07), H, TEAL)

    add_tb(sl, Inches(6.65), Inches(0.3), Inches(6.4), Inches(0.35), [LABEL(11, 'STRATEGY')])
    add_tb(sl, Inches(6.65), Inches(0.7), Inches(6.4), Inches(1.2), [
        TITLE('Technology + Philosophy', size=34),
    ])
    add_hr(sl, Inches(6.65), Inches(1.95), Inches(2.5), TEAL)
    add_tb(sl, Inches(6.65), Inches(2.15), Inches(6.3), Inches(2.3), [
        HEAD('Technology Layer — Living Works Platform'),
        SPC(),
        BUL('SaaS platform for biological design'),
        BUL('Enterprise partnerships and custom design services'),
        BUL('IP licensing across materials, construction, biotech'),
    ])
    add_tb(sl, Inches(6.65), Inches(4.3), Inches(6.3), Inches(2.3), [
        HEAD('Cultural Layer — Spiral Steward Philosophy'),
        SPC(),
        BUL('Sagent Creed — philosophical storytelling'),
        BUL('True Republic — podcast on biology and decentralisation'),
        BUL('Eight Pillars — educational video series'),
        SPC(),
        QTE('"Philosophy creates cultural demand.\nTechnology creates real-world solutions."'),
    ])
    add_tb(sl, Inches(0.2), H - Inches(0.48), Inches(6.0), Inches(0.4), [
        CAP('Fibonacci Spiral Room Layout — designing with living principles'),
    ])
    return sl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
def slide_13_close():
    sl = new_slide(); set_bg(sl, BG)
    add_img(sl, IMGS["dome2"], 0, 0, W, H)
    add_rect(sl, 0, 0, W, H, BG, alpha=45)
    add_rect(sl, 0, 0, Inches(0.07), H, TEAL)

    add_hr(sl, Inches(4.5), Inches(2.8), Inches(4.33), TEAL)
    add_tb(sl, Inches(1.5), Inches(2.9), Inches(10.33), Inches(2.5), [
        L('"Industrial civilisation built with blocks.',
          size=26, italic=True, color=CREAM, align=PP_ALIGN.CENTER, light=True),
        L('The next civilisation will build with living systems."',
          size=26, italic=True, color=CREAM, align=PP_ALIGN.CENTER, light=True, sp=3),
    ])
    add_hr(sl, Inches(4.5), Inches(5.35), Inches(4.33), TEAL)
    add_tb(sl, Inches(1.5), Inches(5.55), Inches(10.33), Inches(0.9), [
        L('LIVING WORKS  ·  Living Works by the Word',
          size=20, bold=True, color=LTGRN, align=PP_ALIGN.CENTER),
    ])
    add_tb(sl, Inches(1.5), Inches(6.4), Inches(10.33), Inches(0.7), [
        L('Living Works provides the tools.  Spiral Stewards guide their use.',
          size=16, color=GREY, align=PP_ALIGN.CENTER, italic=True),
    ])
    return sl

# ── Build slides and attach speaker notes ────────────────────────────────────

add_notes(slide_01_cover(), """\
Good [morning/afternoon] everyone. My name is [your name], and I'm the founder of Living Works by the Word.

What you're looking at is a living structure — a building grown from root systems and mycelium, using the same branching geometry that trees have spent 400 million years perfecting.

This presentation is about a simple but radical idea: that we are entering a new era of civilisation — one that builds with life instead of against it. I call it the Living Age.

Over the next few minutes I'm going to show you what that means technically, why the timing is right, and how we are building the tools to make it happen.
""")

add_notes(slide_02_problem(), """\
Let's start with the problem.

Every building you've ever been in was made the same way: dig up raw materials, burn massive amounts of energy to process them, assemble rigid components, and hope they last long enough before they need to be torn down and replaced.

That's how we've built civilisation for the past two hundred years — and it costs us enormously. Construction accounts for nearly 40% of global carbon emissions. Our materials degrade, they don't adapt, and they have no relationship with the living systems around them.

The left side of this image says it all: The Block. Rigid. Dead matter. Control. Conformity.

We built a world of blocks. And now we're starting to understand the cost of that choice.
""")

add_notes(slide_03_opportunity(), """\
Here's what's remarkable: nature already solved most of the engineering problems we're struggling with.

This image is a mycelium network — the root system of a fungus. It distributes resources across a forest floor with efficiency that would make any logistics engineer envious. No central controller. No wasted energy. Pure adaptive intelligence.

Trees self-assemble structures that handle wind loads, snow loads, and decades of growth. Coral builds mineral architecture that has lasted thousands of years. Our own cells organise into organs without any blueprint telling each one what to do.

The problem isn't that biology can't build. The problem is that our tools only let us observe biology — we have no way to design with it. There's no design software for living systems. That's the gap we're filling.
""")

add_notes(slide_04_vision(), """\
So what's the vision?

This infographic puts it clearly. On the left you have Block World — the industrial civilisation we built. Centralised control, rigid structures, mechanical manufacturing, resource extraction, top-down design.

On the right is what we're moving toward: Spiral World. Decentralised systems. Adaptive growth. Living materials. Regenerative ecosystems. Emergent design.

This is not science fiction. The materials, the biology, the AI tools — they are all converging right now.

Our vision is to give designers, architects, engineers, and scientists the software they need to work in that right-hand world. To make designing with living systems as natural as designing with steel and concrete.

We call this transition the Living Age.
""")

add_notes(slide_05_platform(), """\
This is Living Works — our platform.

Three layers working together.

The media and culture layer — which includes the Sagent Creed, the True Republic podcast, and the Eight Pillars video series — builds the narrative and the community. It answers the question: why should the world care about designing with biology?

The technology layer is the core product: a morphogenesis engine, an AI biological design interface, and a BioCAD platform. These are the actual tools that let you design with life.

And the applications layer shows where this goes in the real world: living architecture, living materials, ecosystem engineering, biological manufacturing.

The next few slides walk through each technology component.
""")

add_notes(slide_06_morpho(), """\
The first and most foundational technology is the Morphogenesis Engine.

What you're seeing here is a 3D simulation of a mycelium dome — a structure whose shape is entirely determined by biological growth rules rather than by a human designer forcing it into a predetermined form.

The morphogenesis engine simulates how living structures develop. You give it starting conditions — nutrient gradients, environmental forces, growth rules — and it shows you what emerges. It predicts branching, cellular division, structural load distribution, developmental pathways.

Think of it as a physics engine, but instead of simulating gravity and collisions, it simulates life.

This kind of tool doesn't exist anywhere right now. The closest things are academic research simulations that aren't connected to any design workflow. We're building the design-ready version.
""")

add_notes(slide_07_ai(), """\
The second component is the AI Biological Design Interface.

What you're watching on the left is one of our early morphogenetic growth simulations — cells dividing, differentiating, organising themselves from a set of simple rules.

The interface sits on top of this engine and makes it accessible. Instead of writing code or configuring parameters, you describe what you want to grow.

You type something like: "Design a branching structure optimised for wind resistance and water collection." The system interprets that as a biological design problem, generates growth models, suggests relevant gene pathways, and runs structural simulations — all from natural language.

This is the same shift that happened in visual AI — where you went from needing to be a programmer to just describing what you want. We're doing that for biology.
""")

add_notes(slide_08_biocad(), """\
The third component is the BioCAD platform — biological computer-aided design.

This is where Living Works becomes a professional design environment. Think of it like AutoCAD or Revit, but instead of drawing walls and beams, you're planning gene edits, designing tissue architectures, and simulating how a living material will behave in a given environment.

The branching tree structure you see here is a real architectural drawing — a living elevation where the structural columns are grown trees, and the canopy is the roof. That drawing was generated using biological growth principles.

BioCAD connects to our Living Materials Library — a database of biological building blocks: mycelium species, plant tissue scaffolds, bacterial cellulose systems, coral-like mineralisation processes. Each entry comes with its growth rules, environmental requirements, and structural properties.

This is the tool that architects, materials scientists, and bioengineers will use to do their daily work.
""")

add_notes(slide_09_arch(), """\
This is what it looks like when you actually build with it.

These structures are grown — not poured, not welded, not assembled. They emerge from biological growth processes guided by design intent. Each dome follows a spiral geometry derived from natural growth mathematics. The roots anchor them. The walls breathe.

I want you to sit with this image for a moment, because this is where we're going.

We are not talking about putting a few plants on a building facade and calling it sustainable. We are talking about buildings that are fundamentally alive — that grow themselves, repair themselves, absorb carbon from the atmosphere, and become more integrated with their environment over time, not less.

This is living architecture. And it is the first major application market for everything we are building.
""")

add_notes(slide_10_applications(), """\
Living architecture is just the first of four major industries this platform addresses.

The second is carbon-negative construction. Buildings grown from biological materials don't just have a low carbon footprint — they actively absorb CO₂ as they grow. A mycelium structure sequesters carbon. An engineered plant facade keeps photosynthesising for the life of the building.

Third is ecosystem engineering — using biological design tools to create landscapes that regenerate soil, purify water, and restore biodiversity. Not passive rewilding, but actively designed living systems.

And fourth is biological manufacturing — growing the materials that today we make from petroleum. Textiles, packaging, structural composites, insulation. Companies like Ecovative and MycoWorks have proven the concept. We're building the design software layer that makes scaling those materials possible.

Each of these is a multi-billion dollar industry in transition. Living Works positions as the design platform that enables all of them.
""")

add_notes(slide_11_biz(), """\
The business model has three revenue streams, and they're designed to build on each other.

The first is our SaaS platform. Subscription access to the morphogenesis engine, AI design interface, and BioCAD tools. We start with universities and synthetic biology labs — they're the early adopters who validate the technology. Then we move into startups and enterprise clients: architecture firms, materials companies, biotech.

The second stream is AI design services — custom projects where we take a client's design challenge and use our platform to generate biological solutions. This generates revenue immediately while also producing data that trains our AI models.

The third is IP licensing. As we develop novel growth algorithms and engineered biological systems, those become intellectual property that other companies license to build their own products.

The roadmap is phased: Year one and two, establish credibility with university pilots. Year three to five, enterprise adoption and architecture partnerships. Year five to ten, living infrastructure at scale.

The floor plan behind me is a house designed on Fibonacci spiral geometry — every room proportion follows the golden ratio. That's what designing with living principles looks like even at the level of spatial planning.
""")

add_notes(slide_12_strategy(), """\
One of the most important things I want you to understand about this initiative is the two-track strategy.

Track one is the technology company. Living Works is a serious biological design platform. It has a clear SaaS revenue model, a defined market, and technology differentiation that is hard to replicate. When we talk to investors and enterprise clients, this is what we lead with.

Track two is the cultural movement. The Sagent Creed, the True Republic podcast, the Eight Pillars series. This is philosophical storytelling — the narrative that explains why the Living Age matters and builds an audience of people who want it to exist.

Why does the movement matter for the business? Because great technology companies aren't just built on good code. They're built on a worldview that people want to belong to. Tesla didn't just sell cars — it sold a vision of the future. Apple didn't just sell computers — it sold a philosophy of design.

The spiral on this slide is a room arrangement designed using sunflower mathematics — 137.5 degrees between each space, the same angle a sunflower uses to pack seeds. It's beautiful because it follows a living principle.

That's what we're doing. Building technology that follows living principles, and telling the story of why that matters.
""")

add_notes(slide_13_close(), """\
I'll leave you with this.

Industrial civilisation built with blocks. Rigid, centralised, mechanical. It gave us remarkable things. But it also gave us a planet in crisis, ecosystems collapsing, and cities that feel fundamentally disconnected from life.

The next civilisation will build differently. It will grow its structures, cultivate its materials, and design its environments the way nature has always designed — through emergence, adaptation, and cooperation with living systems.

We are at the very beginning of that transition. The biology is ready. The AI tools are ready. What's been missing is the design platform — the software that lets architects, engineers, and scientists work at the intersection of technology and life.

That's what Living Works is.

We're building the tools. And we're looking for the people who want to help guide how they're used.

Thank you.
""")

prs.save(OUT)
print(f"✓  Saved → {OUT}")
print(f"   13 slides with speaker notes")
